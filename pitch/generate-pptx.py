#!/usr/bin/env python3
"""Generate Intent Revenue sales pitch PowerPoint from slide content."""

from pathlib import Path

import fitz
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT = ROOT / "Intent-Revenue-Sales-Pitch.pptx"
LOGO_SVG = ROOT.parent / "Intent Logo i1.svg"
LOGO_PNG = ROOT / "assets" / "intent-logo.png"

BG = RGBColor(0, 0, 0)
FG = RGBColor(248, 250, 252)
ACCENT = RGBColor(34, 211, 238)
MUTED = RGBColor(100, 116, 139)
BORDER = RGBColor(30, 41, 59)

SLIDES = [
    {
        "kicker": "INTENT REVENUE",
        "title": "We Grow Revenue.\nBy A Lot.",
        "subtitle": "Growth Partner For Contractors And The Trades",
        "footer": "intentrev.net",
    },
    {
        "kicker": "THE MARKET TODAY",
        "title": "Legacy operators still win on inertia, not quality",
        "bullets": [
            "Slow follow-up and bloated overhead in established shops",
            "Local markets locked by operators who stopped improving",
            "Homeowners choose who shows up on Google, not who does the best work",
            "New operators have the skill but not the systems to take share",
        ],
    },
    {
        "kicker": "OUR MISSION",
        "title": "Turning Small Businesses to Leading Competitors",
        "body": "Intent Revenue gives contractors the revenue systems, software, and search presence that legacy shops never built.",
    },
    {
        "kicker": "BUILT FOR THE TRADES",
        "title": "Revenue strategy and software first",
        "bullets": [
            "#1 Revenue Streams: demand, conversion, repeat work",
            "#2 Application & Software: custom site, intake, dashboard",
            "Google Search: local SEO, geo pages, content",
            "Paid Ads & Content: campaigns that amplify organic",
        ],
    },
    {
        "kicker": "PHONE-FIRST",
        "title": "Trades run on the phone. We build around that.",
        "bullets": [
            "Speed-to-lead and intake automation on every missed call and form",
            "Custom software integrated with CRM, email, and scheduling tools",
            "Dashboards tie marketing spend to leads and booked jobs by source",
            "Organic search first; paid when it accelerates what's working",
        ],
    },
    {
        "kicker": "THE INTENT REVENUE PACKAGE",
        "title": "What's included (1/2)",
        "bullets": [
            "Partnership qualification before we start",
            "Revenue stream discovery and growth",
            "Custom Intent Revenue application and software",
            "Speed-to-lead and intake automation",
            "Google Reviews and Google Business Profile engine",
        ],
    },
    {
        "kicker": "THE INTENT REVENUE PACKAGE",
        "title": "What's included (2/2)",
        "bullets": [
            "Google Search organic growth",
            "Paid ads and content creation",
            "Deep business integration (CRM, email, scheduling, ad accounts)",
            "Ongoing optimization and support",
        ],
    },
    {
        "kicker": "SELECTIVE PARTNERSHIP",
        "title": "Clients we can guarantee results for",
        "bullets": [
            "5+ jobs/mo or $25k+ annual revenue",
            "Phone-first operations",
            "Google Reviews on active Google Business Profile",
            "Defined territory",
            "Growth investment",
            "Full business access (CRM, email, scheduling, ad accounts)",
            "Owner in the fight",
        ],
    },
    {
        "kicker": "PATH TO QUALIFICATION",
        "title": "Intent Launchpad",
        "bullets": [
            "Post-job Google review engine",
            "Google Business Profile build-out",
            "Intake basics: missed-call text-back, after-hours, lead form",
            "Monthly milestones until full partnership",
        ],
    },
    {
        "kicker": "REVENUE ENGINEERING",
        "title": "Not vanity metrics. Booked jobs.",
        "bullets": [
            "Leads, calls, and outcomes tracked by source",
            "See what's paying for itself before you scale spend",
            "Built for your trade, not a generic marketing dashboard",
        ],
    },
    {
        "kicker": "NEXT STEP",
        "title": "Two ways to work with Intent Revenue",
        "bullets": [
            "Qualified: Full Partnership — revenue, software, search, and ads",
            "Not there yet: Launchpad — reviews, intake, milestones first",
        ],
    },
    {
        "kicker": "LET'S GO",
        "title": "Ready to grow revenue, not just traffic?",
        "body": "Tell us your trade and territory. Intent Revenue will assess fit and map what we build for your business.",
        "footer": "intentrev.net · Apply for Partnership · Start with Launchpad",
    },
]


def ensure_logo_png() -> Path:
    """Rasterize Intent Logo i1.svg for PowerPoint (PPTX does not embed SVG)."""
    LOGO_PNG.parent.mkdir(parents=True, exist_ok=True)
    if not LOGO_SVG.exists():
        raise FileNotFoundError(f"Logo not found: {LOGO_SVG}")

    svg_mtime = LOGO_SVG.stat().st_mtime
    png_mtime = LOGO_PNG.stat().st_mtime if LOGO_PNG.exists() else 0
    if not LOGO_PNG.exists() or svg_mtime > png_mtime:
        doc = fitz.open(str(LOGO_SVG))
        page = doc[0]
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=True)
        pix.save(str(LOGO_PNG))

    return LOGO_PNG


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def style_run(run, *, size=18, color=FG, bold=False, font_name="Segoe UI"):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name


def add_navbar(slide, prs, logo_path: Path, *, is_title: bool):
    """Fixed top bar with Intent logo — large on slide 1, compact on all others."""
    nav_h = Inches(1.05) if is_title else Inches(0.55)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        nav_h,
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = BG
    bar.line.color.rgb = BORDER
    bar.line.width = Pt(1)

    logo_h = Inches(0.78) if is_title else Inches(0.28)
    logo_top = Inches(0.12) if is_title else Inches(0.13)
    slide.shapes.add_picture(str(logo_path), Inches(0.35), logo_top, height=logo_h)

    if is_title:
        brand = add_textbox(slide, Inches(1.35), Inches(0.2), Inches(5.5), Inches(0.45))
        tf = brand.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Intent Revenue"
        style_run(run, size=22, bold=True)

        tag = add_textbox(slide, Inches(1.35), Inches(0.58), Inches(6), Inches(0.3))
        tf = tag.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "WE GROW REVENUE. BY A LOT."
        style_run(run, size=8, color=ACCENT, font_name="Consolas")


def add_kicker(slide, text, top):
    box = add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.4))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text.upper()
    style_run(run, size=10, color=ACCENT, font_name="Consolas")
    p.alignment = PP_ALIGN.LEFT


def add_title(slide, text, top):
    box = add_textbox(slide, Inches(0.6), top, Inches(12.1), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        style_run(run, size=34 if i == 0 else 30, bold=True)
    return box


def add_body(slide, text, top):
    box = add_textbox(slide, Inches(0.6), top, Inches(12.1), Inches(1.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    style_run(run, size=18, color=FG)
    return box


def add_bullets(slide, items, top):
    box = add_textbox(slide, Inches(0.6), top, Inches(12.1), Inches(4.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = item
        style_run(run, size=17)
    return box


def add_footer(slide, text):
    box = add_textbox(slide, Inches(0.6), Inches(6.85), Inches(12), Inches(0.35))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    style_run(run, size=10, color=MUTED, font_name="Consolas")


def add_subtitle(slide, text, top):
    box = add_textbox(slide, Inches(0.6), top, Inches(12), Inches(0.5))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    style_run(run, size=20, color=MUTED)


def layout_tops(is_title: bool):
    if is_title:
        return {
            "kicker": Inches(1.35),
            "title": Inches(1.9),
            "subtitle": Inches(3.25),
            "body": Inches(3.4),
            "bullets": Inches(3.15),
        }
    return {
        "kicker": Inches(0.72),
        "title": Inches(1.22),
        "subtitle": Inches(2.32),
        "body": Inches(2.45),
        "bullets": Inches(2.3),
    }


def build():
    logo_path = ensure_logo_png()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for i, data in enumerate(SLIDES):
        is_title = i == 0
        tops = layout_tops(is_title)

        slide = prs.slides.add_slide(blank)
        set_slide_bg(slide)
        add_navbar(slide, prs, logo_path, is_title=is_title)

        if data.get("kicker"):
            add_kicker(slide, data["kicker"], tops["kicker"])
        if data.get("title"):
            add_title(slide, data["title"], tops["title"])
        if data.get("subtitle"):
            add_subtitle(slide, data["subtitle"], tops["subtitle"])
        if data.get("body"):
            add_body(slide, data["body"], tops["body"])
        if data.get("bullets"):
            add_bullets(slide, data["bullets"], tops["bullets"])
        if data.get("footer"):
            add_footer(slide, data["footer"])

    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
